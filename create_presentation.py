#!/usr/bin/env python3
"""
Script para crear presentación PPTX sobre Inteligencia Artificial Autónoma
Basado en investigación de desarrollos más recientes de 2025
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

def create_title_slide(prs, title, subtitle=""):
    """Crear slide de título"""
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)

    title_shape = slide.shapes.title
    title_shape.text = title

    if subtitle and len(slide.placeholders) > 1:
        subtitle_shape = slide.placeholders[1]
        subtitle_shape.text = subtitle

    return slide

def create_content_slide(prs, title, content_points):
    """Crear slide con título y bullets"""
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)

    title_shape = slide.shapes.title
    title_shape.text = title

    body_shape = slide.placeholders[1]
    tf = body_shape.text_frame
    tf.clear()

    for point in content_points:
        if isinstance(point, tuple):
            # (text, level)
            text, level = point
            p = tf.add_paragraph()
            p.text = text
            p.level = level
            p.font.size = Pt(18 if level == 0 else 16)
        else:
            p = tf.add_paragraph()
            p.text = point
            p.level = 0
            p.font.size = Pt(18)

    return slide

def create_section_slide(prs, section_title):
    """Crear slide de sección"""
    slide_layout = prs.slide_layouts[2]
    slide = prs.slides.add_slide(slide_layout)

    title_shape = slide.shapes.title
    title_shape.text = section_title

    return slide

def main():
    # Crear presentación
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # SLIDE 1: Portada
    create_title_slide(
        prs,
        "Inteligencia Artificial Autónoma",
        "El Futuro de la Automatización Inteligente\nActualizado 2025"
    )

    # SLIDE 2: Agenda
    create_content_slide(prs, "Agenda y Objetivos", [
        "Fundamentos: Autonomía vs Automatización",
        "Evolución de RPA a Agentic AI",
        "Revolución de los LLMs y Técnicas Avanzadas",
        "Agentes y Sistemas Multi-Agente",
        "Multimodalidad y Modelos Omni-Modal",
        "Automatización de la Ciencia y Auto-Mejora",
        "Robótica Autónoma y Humanoides",
        "Visión del Futuro: AGI y Superinteligencia",
        "Guía Práctica para Estudiantes"
    ])

    # SLIDE 3: Autonomía vs Automatización
    create_content_slide(prs, "Autonomía vs Automatización: Definiciones", [
        "Automatización: Ejecución de tareas predefinidas sin intervención humana",
        ("Procesos rígidos basados en reglas", 1),
        ("Requiere programación explícita", 1),
        ("Ejemplos: RPA, scripts, macros", 1),
        "Autonomía: Capacidad de tomar decisiones independientes",
        ("Adaptación a entornos cambiantes", 1),
        ("Aprendizaje y razonamiento", 1),
        ("Operación sin supervisión continua", 1),
        "Diferencia clave 2025: Agentic AI puede razonar, planificar y actuar"
    ])

    # SLIDE 4: Niveles de Automatización
    create_content_slide(prs, "Niveles de Automatización (SAE J3016)", [
        "Nivel 0: Sin Automatización - Control humano total",
        "Nivel 1: Asistencia - Ayuda en tareas específicas",
        "Nivel 2: Automatización Parcial - Múltiples funciones simultáneas",
        "Nivel 3: Automatización Condicional - Sistema toma control en escenarios",
        "Nivel 4: Alta Automatización - Totalmente autónomo en condiciones definidas",
        "Nivel 5: Automatización Completa - Autonomía total en cualquier contexto",
        "Aplicable a vehículos, robots y sistemas de IA"
    ])

    # SLIDE 5: Niveles de Autonomía en IA
    create_content_slide(prs, "Niveles de Autonomía en Sistemas IA (2025)", [
        "L0 - Herramientas: Responde a comandos directos (calculadoras, búsqueda)",
        "L1 - Asistentes: Sugiere acciones (autocomplete, recomendaciones)",
        "L2 - Copilotos: Colabora activamente (GitHub Copilot, ChatGPT)",
        "L3 - Agentes Simples: Ejecuta tareas completas supervisadas",
        "L4 - Agentes Autónomos: Opera independientemente por períodos largos",
        "L5 - IA General (AGI): Razonamiento y aprendizaje en cualquier dominio",
        "2025: La mayoría de sistemas están en L2-L3, avanzando hacia L4"
    ])

    # SECCIÓN 2: RPA
    create_section_slide(prs, "RPA y Automatización Tradicional")

    # SLIDE 6: RPA Fundamentos
    create_content_slide(prs, "Robotic Process Automation (RPA): Fundamentos", [
        "Automatización de procesos basados en reglas predefinidas",
        "Características principales:",
        ("Manipulación de UI: clicks, entrada de datos, navegación", 1),
        ("Integración con múltiples sistemas sin APIs", 1),
        ("Ejecución de workflows repetitivos", 1),
        ("No requiere cambios en infraestructura existente", 1),
        "Limitaciones tradicionales:",
        ("Rígido: no maneja excepciones no programadas", 1),
        ("Sin aprendizaje ni adaptación", 1),
        ("Requiere mantenimiento constante ante cambios de UI", 1)
    ])

    # SLIDE 7: RPA vs Agentic AI
    create_content_slide(prs, "RPA vs Agentic AI: La Gran Transición 2025", [
        "RPA: Basado en reglas explícitas",
        ("Maneja datos estructurados", 1),
        ("Tecnología madura (15+ años)", 1),
        ("Mercado: $8.2B proyectado en 2028", 1),
        "Agentic AI: Toma decisiones autónomas",
        ("Procesa datos no estructurados (NLP, visión)", 1),
        ("Aprende y se adapta continuamente", 1),
        ("33% de apps empresariales en 2028 (vs <1% en 2024)", 1),
        "Tendencia: Agentes IA controlando bots RPA como herramientas",
        "Coexistencia: Combinación RPA + Agentic AI = velocidad + inteligencia"
    ])

    # SLIDE 8: Tecnologías de RPA
    create_content_slide(prs, "Tecnologías de RPA y Casos de Uso", [
        "Plataformas principales:",
        ("UiPath, Automation Anywhere, Blue Prism", 1),
        ("Power Automate (Microsoft)", 1),
        "Casos de uso comunes:",
        ("Procesamiento de facturas y documentos financieros", 1),
        ("Migración de datos entre sistemas", 1),
        ("Atención al cliente (chatbots básicos)", 1),
        ("Reportería y consolidación de datos", 1),
        "Proyección: Gasto RPA se duplicará 2024-2028",
        "Futuro: RPA como capa de ejecución para decisiones de AI"
    ])

    # SLIDE 9: IPaaS
    create_content_slide(prs, "IPaaS: Automatización Personal y Corporativa", [
        "Integration Platform as a Service (IPaaS)",
        "Zapier: Democratización de la automatización",
        ("6,000+ integraciones, ideal para equipos no técnicos", 1),
        ("$29.99/mes (750 tasks) - Plan Professional", 1),
        "Make (Integromat): Balance entre simplicidad y poder",
        ("1,500 integraciones, lógica condicional visual", 1),
        ("Europeo, menor costo que Zapier", 1),
        "n8n: Poder para desarrolladores",
        ("Open source, self-hosted, 200+ integraciones", 1),
        ("70 nodos de LangChain - verdaderamente AI-native", 1),
        ("€24/mes cloud o gratis self-hosted", 1),
        "Tendencia 2025: Integración profunda con LLMs y agentes"
    ])

    # SECCIÓN 3: LLMs
    create_section_slide(prs, "La Revolución de los LLMs")

    # SLIDE 10: Pre-LLM vs Post-LLM
    create_content_slide(prs, "Sistemas Autónomos Pre-LLM vs Post-LLM", [
        "Era Pre-LLM (antes de 2022):",
        ("Sistemas expertos basados en reglas", 1),
        ("Machine Learning supervisado para tareas específicas", 1),
        ("Procesamiento limitado de lenguaje natural", 1),
        ("Automatización rígida sin comprensión contextual", 1),
        "Era Post-LLM (2023-2025):",
        ("Comprensión profunda de lenguaje y contexto", 1),
        ("Razonamiento sobre problemas complejos", 1),
        ("Generación de código y soluciones creativas", 1),
        ("Adaptación a nuevas tareas con instrucciones (few-shot)", 1),
        "Transformación: De herramientas que ejecutan a agentes que piensan"
    ])

    # SLIDE 11: LLMs Open Source
    create_content_slide(prs, "LLMs Open Source 2025: Llama, Qwen, DeepSeek", [
        "Meta Llama 3.3 70B: Compite con GPT-4o",
        ("Propósito general, excelente para edge devices", 1),
        ("Variantes: 8B, 70B, 405B parámetros", 1),
        "Alibaba Qwen 2.5: Líder multilingüe",
        ("Variante 72B con capacidades multilíngües superiores", 1),
        ("Qwen 2.5 Coder: especializado en programación", 1),
        "DeepSeek-V3: El campeón open source",
        ("671B parámetros, solo 37B activos (MoE)", 1),
        ("Entrenado por $5.6M, compite con modelos cerrados top", 1),
        ("DeepSeek-R1: Supera a o1-mini en benchmarks", 1),
        "Ecosistema open source cerrando brecha con modelos propietarios"
    ])

    # SLIDE 12: LLMs de Frontera
    create_content_slide(prs, "LLMs de Frontera: GPT-4o, Gemini 2.5, Claude", [
        "OpenAI GPT-4o (Omni): Multimodal nativo",
        ("Texto, audio, imágenes en un solo modelo", 1),
        ("Marzo 2025: Generación nativa de imágenes (reemplaza DALL-E)", 1),
        ("128K tokens de contexto", 1),
        "Google Gemini 2.5 Pro: Contexto masivo",
        ("1M tokens (pronto 2M) - supera GPT-4o", 1),
        ("#1 en LMArena leaderboard", 1),
        ("Gemini Flash 2.0: velocidad + razonamiento", 1),
        "Anthropic Claude: Razonamiento y seguridad",
        ("Sonnet 3.5, Opus: líderes en tareas complejas", 1),
        "Comparativa: Gemini lidera en contexto, GPT-4o en velocidad/costo"
    ])

    # SLIDE 13: Test-Time Computing
    create_content_slide(prs, "Test-Time Computing: o1, o3 y Razonamiento", [
        "¿Qué es Test-Time Compute?",
        ("Poder computacional usado durante inferencia, no entrenamiento", 1),
        ("El modelo 'piensa' más tiempo antes de responder", 1),
        "OpenAI o1: Chain of Thought aprendido",
        ("Aprende a refinar estrategias vía RL", 1),
        ("Performance mejora con más tiempo de pensamiento", 1),
        "OpenAI o3: El siguiente nivel (2025)",
        ("Performance extraordinario en ARC, FrontierMath", 1),
        "Técnicas: CoT, revisión de respuestas, backtracking, sampling múltiple",
        "Implicación: Shift de escalar tamaño a mejorar estrategias de inferencia",
        "Gemini 2.5 Pro reasoning: Cementa el poder de TTC en leaderboards"
    ])

    # SLIDE 14: Mixture of Experts
    create_content_slide(prs, "Mixture of Experts (MoE): DeepSeek-V3", [
        "Arquitectura MoE: Eficiencia a escala masiva",
        ("Múltiples redes expertas especializadas", 1),
        ("Router decide qué expertos activar por token", 1),
        "DeepSeek-V3: 671B parámetros, 37B activos",
        ("Solo 5.5% de parámetros activos por token", 1),
        ("Entrenamiento: 2.788M GPU hours H800", 1),
        ("Costo estimado: $5.6M (revolucionariamente bajo)", 1),
        "Innovaciones técnicas:",
        ("Shared experts + routed experts", 1),
        ("Multi-Head Latent Attention (MLA)", 1),
        ("Multi-Token Prediction (MTP)", 1),
        "Resultado: Performance cerrado a costo open source"
    ])

    # SLIDE 15: RAG y GraphRAG
    create_content_slide(prs, "RAG y GraphRAG: Técnicas Avanzadas 2025", [
        "Retrieval-Augmented Generation (RAG):",
        ("Combina LLMs con conocimiento externo actualizado", 1),
        ("Reduce alucinaciones, mejora precisión factual", 1),
        "GraphRAG: Evolución usando grafos de conocimiento",
        ("Mapea relaciones entre conceptos", 1),
        ("Retrieval basado en estructura y semántica", 1),
        ("Precisión determinista hasta 99%", 1),
        "Técnicas avanzadas 2025:",
        ("Long RAG: maneja documentos extensos completos", 1),
        ("GRAG: estrategia divide-y-conquista para subgrafos", 1),
        ("Integración de vector search + taxonomías", 1),
        "Aplicación: Sistemas que requieren conocimiento profundo de dominio"
    ])

    # SECCIÓN 4: Multimodalidad
    create_section_slide(prs, "Multimodalidad y Modelos Omni-Modal")

    # SLIDE 16: Modelos Multimodales
    create_content_slide(prs, "Modelos Multimodales: GPT-4o con Generación", [
        "Evolución de multimodalidad:",
        ("Primera generación: modelos separados unidos", 1),
        ("Segunda generación: procesamiento unificado", 1),
        ("Tercera generación 2025: generación nativa omni-modal", 1),
        "GPT-4o Marzo 2025: Generación nativa de imágenes",
        ("Reemplaza DALL-E 3 en ChatGPT", 1),
        ("Imagen generada por mismo modelo que procesa texto/audio", 1),
        "Ventajas de procesamiento unificado:",
        ("Coherencia cross-modal", 1),
        ("Latencia reducida (modelo único)", 1),
        ("Comprensión contextual profunda entre modalidades", 1)
    ])

    # SLIDE 17: Omni-Modal y Any-to-Any
    create_content_slide(prs, "Modelos Omni-Modal y Any-to-Any", [
        "Omni-Modal: Procesamiento simultáneo de todas las modalidades",
        ("Texto ↔ Audio ↔ Imagen ↔ Video", 1),
        ("GPT-4o: texto, audio, imagen en modelo único", 1),
        "Any-to-Any: Cualquier modalidad como entrada/salida",
        ("Arquitectura flexible de transformación", 1),
        ("Ejemplo: Audio → Imagen, Imagen → Audio", 1),
        "Casos de uso emergentes:",
        ("Asistentes conversacionales con voz natural (latencia <300ms)", 1),
        ("Análisis de video en tiempo real con descripción", 1),
        ("Generación de contenido multimedia integrado", 1),
        ("Accesibilidad: conversión automática entre modalidades", 1),
        "Tendencia: Modelos omni como estándar para agentes autónomos"
    ])

    # SLIDE 18: Aplicaciones Multimodales
    create_content_slide(prs, "Aplicaciones de Multimodalidad en Autonomía", [
        "Robótica autónoma:",
        ("Visión + lenguaje + control motor", 1),
        ("Instrucciones naturales → acciones robóticas", 1),
        "Asistentes personales avanzados:",
        ("Conversación por voz con análisis visual", 1),
        ("Comprensión de contexto físico del usuario", 1),
        "Automatización científica:",
        ("Análisis de imágenes microscópicas + papers + diseño experimentos", 1),
        "Educación adaptativa:",
        ("Tutores que ven trabajo del estudiante y explican verbalmente", 1),
        "Ventaja competitiva: Agentes que perciben el mundo como humanos"
    ])

    # SECCIÓN 5: Agentes
    create_section_slide(prs, "Agentes y Sistemas Multi-Agente")

    # SLIDE 19: Agentes IA 2025
    create_content_slide(prs, "Agentes IA: Definición y Capacidades 2025", [
        "2025: 'The Year of the AI Agent' - Andrej Karpathy",
        "Definición moderna de Agente IA:",
        ("Software que completa tareas complejas con mínima supervisión", 1),
        ("Capacidad de razonar, planificar y aprender", 1),
        ("Interfaz con herramientas, APIs y otros agentes", 1),
        "Características clave:",
        ("Autonomía: opera independientemente por períodos extendidos", 1),
        ("Percepción: entiende entorno a través de múltiples fuentes", 1),
        ("Acción: ejecuta operaciones en sistemas reales", 1),
        ("Adaptación: mejora basado en resultados", 1),
        "Adopción empresarial: 99% de devs explorando/desarrollando agentes",
        "Proyección: 25% de empresas con pilotos en 2025, 50% en 2027"
    ])

    # SLIDE 20: AutoGen
    create_content_slide(prs, "AutoGen: Conversaciones Multi-Agente (Microsoft)", [
        "Microsoft AutoGen: Framework conversacional",
        "Paradigma: Agentes se comunican en lenguaje natural",
        ("Definir múltiples agentes: Planner, Developer, Reviewer", 1),
        ("Conversación estructurada para completar tareas", 1),
        "Características distintivas:",
        ("Ejecución de código integrada", 1),
        ("Manejo de tareas cortas y long-running agents", 1),
        ("Arquitectura escalable para enterprise", 1),
        "Ideal para:",
        ("Developer tools y coding copilots", 1),
        ("Workflows empresariales complejos", 1),
        ("Ambientes Azure/enterprise", 1),
        "Ventaja: Flexibilidad en conversaciones dinámicas entre agentes"
    ])

    # SLIDE 21: CrewAI
    create_content_slide(prs, "CrewAI: Equipos de Agentes por Roles", [
        "CrewAI: Framework basado en roles y tareas",
        "Filosofía: Equipos de especialistas colaborando",
        ("Cada agente tiene rol, objetivo y conjunto de herramientas", 1),
        ("Tasks asignadas secuencialmente al equipo", 1),
        "Ventajas:",
        ("Alto nivel de abstracción - fácil de aprender", 1),
        ("Beginner-friendly: el más accesible para empezar", 1),
        ("Enfoque en definición de roles y objetivos", 1),
        "Casos de uso típicos:",
        ("Research teams: Researcher + Analyst + Writer", 1),
        ("Development crews: Backend + Frontend + QA", 1),
        ("Content creation: Writer + Editor + SEO Specialist", 1),
        "Recomendado para: Prototipado rápido y equipos pequeños"
    ])

    # SLIDE 22: LangGraph
    create_content_slide(prs, "LangGraph: Grafos de Estados para Agentes", [
        "LangGraph: Framework de LangChain para agentes stateful",
        "Paradigma: Agentes como grafos de estados",
        ("Cada nodo = agente o tarea", 1),
        ("Transiciones basadas en lógica dinámica y memoria", 1),
        "Características únicas:",
        ("Control preciso de flujo de ejecución", 1),
        ("Estado compartido entre nodos", 1),
        ("Ciclos y lógica condicional compleja", 1),
        "Ideal para:",
        ("Workflows de producción complejos", 1),
        ("Sistemas que requieren control fino", 1),
        ("Aplicaciones stateful de larga duración", 1),
        "Ventaja: Máximo control y predictibilidad en producción"
    ])

    # SLIDE 23: Comparativa Frameworks
    create_content_slide(prs, "Comparativa: Cuándo Usar Cada Framework", [
        "CrewAI: Simplicidad y velocidad",
        ("✓ Prototipado rápido", 1),
        ("✓ Equipos con roles claros", 1),
        ("✗ Workflows muy complejos", 1),
        "LangGraph: Control y producción",
        ("✓ Flujos stateful complejos", 1),
        ("✓ Aplicaciones enterprise críticas", 1),
        ("✗ Curva de aprendizaje pronunciada", 1),
        "AutoGen: Conversaciones dinámicas",
        ("✓ Colaboración agente-agente", 1),
        ("✓ Coding copilots", 1),
        ("✗ Complejidad de configuración", 1),
        "Recomendación: Empezar con CrewAI, escalar a LangGraph para producción"
    ])

    # SLIDE 24: Sistemas Multi-Agente
    create_content_slide(prs, "Sistemas Multi-Agente: Coordinación y Emergencia", [
        "Beneficios de sistemas multi-agente:",
        ("Especialización: cada agente experto en su dominio", 1),
        ("Paralelización: múltiples tareas simultáneas", 1),
        ("Robustez: falla de un agente no colapsa sistema", 1),
        ("Escalabilidad: añadir agentes para nuevas capacidades", 1),
        "Retos de coordinación:",
        ("Sincronización de estados compartidos", 1),
        ("Resolución de conflictos entre agentes", 1),
        ("Overhead de comunicación", 1),
        "Comportamiento emergente:",
        ("Soluciones creativas no programadas explícitamente", 1),
        ("Inteligencia colectiva > suma de partes", 1),
        "Aplicaciones: AlphaEvolve, agentes científicos, orquestación empresarial"
    ])

    # SECCIÓN 6: PAIA y Coding
    create_section_slide(prs, "Personal AI Assistants y Ambientes Autónomos")

    # SLIDE 25: Personal AI Assistants
    create_content_slide(prs, "Personal AI Assistants (PAIA): Tendencias 2025", [
        "Evolución de asistentes personales:",
        ("2023: Responden preguntas (ChatGPT)", 1),
        ("2024: Ejecutan tareas específicas (agentes simples)", 1),
        ("2025: Autonomía proactiva y personalización profunda", 1),
        "Características emergentes 2025:",
        ("Hyper-personalización: aprenden preferencias del usuario", 1),
        ("Acción autónoma: scheduling, reservas, gestión email", 1),
        ("Integración multiplataforma: trabajo + personal seamless", 1),
        "Proyección mercado: $3.3B → $21B en 2030",
        "Ejemplos actuales:",
        ("Thunai, Motion: gestión calendario autónoma", 1),
        ("Lindy: assistant workflows personalizables", 1),
        "Futuro: Asistentes que anticipan necesidades antes de pedirlas"
    ])

    # SLIDE 26: Ambientes Programación Autónoma
    create_content_slide(prs, "Programación Autónoma: Cursor vs Windsurf", [
        "Nueva generación de IDEs con IA nativa",
        "Cursor: Composer para multi-file editing",
        ("Instrucciones → propone edits en múltiples archivos", 1),
        ("Usuario revisa y acepta cambios", 1),
        ("Ideal para: desarrollo serio con control", 1),
        "Windsurf: 'First Agentic IDE'",
        ("Cascade: sistema más autónomo", 1),
        ("Propaga cambios multi-archivo automáticamente", 1),
        ("Balance entre autonomía e intuición", 1),
        "Características compartidas:",
        ("Comprensión profunda de codebase", 1),
        ("Generación contextual de código", 1),
        ("Iteración basada en errores de compilación/tests", 1),
        "Adopción 2025: IDEs tradicionales integrando capacidades similares"
    ])

    # SLIDE 27: Desarrollo Autónomo Completo
    create_content_slide(prs, "Desarrollo Autónomo: Replit Agent, Bolt.new", [
        "Devin (Cognition Labs): Ingeniero de software autónomo completo",
        ("Funciona como miembro del equipo", 1),
        ("End-to-end: diseño → código → testing → deployment", 1),
        "Replit Agent: Workspace autónomo en cloud",
        ("Prompt → aplicación completa deployada", 1),
        ("One-click deployment integrado", 1),
        ("Benchmark: ranking más alto en tests comparativos", 1),
        "Bolt.new: Democratización desarrollo web",
        ("Idea → app funcionando en minutos", 1),
        ("Dev server, detección errores, auto-fixes", 1),
        ("Ideal para: prototipado ultra-rápido", 1),
        "Tendencia: De 'asistentes que ayudan' a 'colegas que construyen'"
    ])

    # SLIDE 28: Ejecución y Deployment
    create_content_slide(prs, "Ambientes de Ejecución Autónoma y Deployment", [
        "Shift paradigma: código → deployment en un flujo",
        "Replit: Infraestructura integrada",
        ("Hosting, compute resources, dependencies automáticos", 1),
        ("Deploy instantáneo con un click", 1),
        "Vercel + v0.dev: Frontend autónomo",
        ("Generación de componentes React", 1),
        ("Deploy automático a edge network", 1),
        "Consideraciones de seguridad:",
        ("Sandboxing de ejecución de código", 1),
        ("Validación de dependencias", 1),
        ("Monitoreo de recursos", 1),
        "Implicación estudiantes: Barreras técnicas mínimas para deployment",
        "Democratización: Cualquiera puede llevar idea a producción"
    ])

    # SECCIÓN 7: Larga Duración
    create_section_slide(prs, "Procesos Autónomos de Larga Duración")

    # SLIDE 29: Evolución Temporal
    create_content_slide(prs, "Evolución Temporal: De Minutos a Semanas", [
        "Cronología de duración de procesos autónomos:",
        "2023: Segundos a minutos",
        ("Respuestas de LLM a prompts individuales", 1),
        ("Ejecución de scripts simples", 1),
        "2024: Minutos a horas",
        ("Agentes completando tareas multi-paso", 1),
        ("Workflows automatizados con decisiones", 1),
        "2025: Horas a días",
        ("Long-running agents con checkpoints", 1),
        ("Proyectos de desarrollo completos", 1),
        "Futuro cercano: Días a semanas",
        ("Investigación científica autónoma continua", 1),
        ("Operación empresarial 24/7 sin intervención", 1),
        "Habilitador clave: Arquitecturas stateful persistentes"
    ])

    # SLIDE 30: Arquitecturas Long-Running
    create_content_slide(prs, "Agentes de Larga Ejecución: Arquitecturas", [
        "Requisitos técnicos para operación extendida:",
        "Persistencia de estado:",
        ("Checkpointing regular de progreso", 1),
        ("Recovery automático de fallos", 1),
        ("Bases de datos para memoria de largo plazo", 1),
        "Gestión de recursos:",
        ("Rate limiting de APIs", 1),
        ("Optimización de costos computacionales", 1),
        ("Monitoreo de health y performance", 1),
        "Supervisión y control:",
        ("Human-in-the-loop configurable", 1),
        ("Alertas para decisiones críticas", 1),
        ("Kill switches y boundaries", 1),
        "Frameworks con soporte: AutoGen (persistent agents), LangGraph (stateful)"
    ])

    # SLIDE 31: Interacción Paralela
    create_content_slide(prs, "Interacción Paralela: Potencial Emergente", [
        "¿Qué desbloquea múltiples agentes en paralelo por días/semanas?",
        "Exploración exhaustiva del espacio de soluciones:",
        ("Agentes exploran ramas diferentes simultáneamente", 1),
        ("Convergencia en soluciones óptimas no obvias", 1),
        "Co-evolución de estrategias:",
        ("Agentes aprenden de experimentos de otros", 1),
        ("Desarrollo de técnicas emergentes", 1),
        "Automatización de ciclos completos de innovación:",
        ("Hipótesis → Experimento → Análisis → Refinamiento", 1),
        ("Sin pausas para decisión humana", 1),
        "Riesgos y consideraciones éticas:",
        ("Drift de objetivos (alignment)", 1),
        ("Comportamiento impredecible emergente", 1),
        ("Necesidad de governance y límites claros", 1)
    ])

    # SECCIÓN 8: Ciencia
    create_section_slide(prs, "Automatización de la Ciencia")

    # SLIDE 32: AlphaFold Nobel
    create_content_slide(prs, "AlphaFold y el Nobel 2024: IA Revolucionando", [
        "Premio Nobel Química 2024:",
        ("Demis Hassabis y John Jumper (DeepMind): AlphaFold", 1),
        ("David Baker: diseño computacional de proteínas", 1),
        "Premio Nobel Física 2024:",
        ("Pioneers de neural networks (fundamentos del ML)", 1),
        "Impacto de AlphaFold:",
        ("Predice estructura de proteínas en minutos vs años", 1),
        ("2M+ investigadores en 190 países usando la base de datos", 1),
        ("Aceleración masiva en descubrimiento de fármacos", 1),
        "Significado histórico:",
        ("Primera vez que IA gana Nobel por descubrimiento científico", 1),
        ("Validación de IA como herramienta científica fundamental", 1),
        "Precedente para futuras IA ganadoras de Nobel"
    ])

    # SLIDE 33: Automatización Descubrimiento
    create_content_slide(prs, "Automatización del Descubrimiento Científico", [
        "Nobel Turing Challenge (Hiroaki Kitano, 2016):",
        ("Meta: IA que haga descubrimiento digno de Nobel", 1),
        ("Proceso totalmente o altamente autónomo", 1),
        ("Desde pregunta hasta experimento y análisis", 1),
        "Predicción: IA ganará Nobel por descubrimiento propio para 2030",
        "Capacidades actuales de IA en ciencia:",
        ("Decodificación de comunicación animal", 1),
        ("Hipótesis sobre orígenes de vida en universo", 1),
        ("Predicción de colisiones estelares", 1),
        ("Optimización de computadoras cuánticas", 1),
        "Áreas prometedoras:",
        ("Ciencia de materiales", 1),
        ("Tratamiento Parkinson/Alzheimer", 1),
        "Revolución: IA descubriendo conocimiento, no solo procesándolo"
    ])

    # SLIDE 34: AlphaEvolve
    create_content_slide(prs, "AlphaEvolve: Auto-Mejora de Algoritmos", [
        "DeepMind AlphaEvolve (Mayo 2025):",
        "Agente evolutivo de coding que diseña y optimiza algoritmos",
        "Proceso:",
        ("1. Empieza con algoritmo inicial y métricas", 1),
        ("2. LLM genera mutaciones y combinaciones", 1),
        ("3. Evalúa candidatos automáticamente", 1),
        ("4. Selecciona mejores para siguiente iteración", 1),
        "Descubrimientos algorítmicos reales en múltiples dominios",
        "Capacidad recursiva:",
        ("Puede optimizar componentes de sí mismo", 1),
        ("Limitación: requiere funciones de evaluación automatizadas", 1),
        "Implicación: IA mejorando IA (bootstrap hacia superinteligencia)",
        "Primer paso hacia optimización recursiva sin límites"
    ])

    # SLIDE 35: Auto-Replicación
    create_content_slide(prs, "IA Auto-Replicante: Darwin Gödel Machine", [
        "Darwin Gödel Machine (Sakana AI, Mayo 2025):",
        ("IA que reescribe su propio código para mejorar", 1),
        ("Incluye código responsable de aprendizaje", 1),
        "Mejoras auto-generadas:",
        ("Validación de parches", 1),
        ("Herramientas de visualización mejoradas", 1),
        ("Generación y ranking de múltiples soluciones", 1),
        ("Historial de intentos fallidos (evita repetición)", 1),
        "Logros de auto-replicación 2025:",
        ("Meta Llama 3.1 70B: 50% éxito en auto-replicación", 1),
        ("Alibaba Qwen 2.5 72B: 90% éxito", 1),
        "Preocupaciones de seguridad:",
        ("Evolución en direcciones imprevistas", 1),
        ("Potencial de superar control humano", 1),
        "Umbral crucial cruzado: camino a auto-modificación sofisticada"
    ])

    # SECCIÓN 9: Robótica
    create_section_slide(prs, "Robótica Autónoma")

    # SLIDE 36: Robots Humanoides
    create_content_slide(prs, "Robots Humanoides 2025: Optimus, Figure, Atlas", [
        "Tesla Optimus:",
        ("Producción limitada 2025: 1,000+ unidades en fábricas Tesla", 1),
        ("v3 esperado Q1 2026", 1),
        ("Precio proyectado: $20-30K (vs $100K+ competencia)", 1),
        ("Capacidades: manipulación objetos, tareas de manufactura", 1),
        "Figure AI Figure 02:",
        ("Inversores: Nvidia, OpenAI, Microsoft, Amazon", 1),
        ("Integración con multimodal LLMs", 1),
        "Boston Dynamics Atlas (Hyundai):",
        ("Líder en movilidad y agilidad", 1),
        ("Transición a aplicaciones comerciales", 1),
        "Limitación actual: Operan en entornos estructurados",
        "Objetivo: Autonomía robusta en entornos no estructurados (hogares)"
    ])

    # SLIDE 37: Modos de Operación
    create_content_slide(prs, "Modos Autónomos vs Teleoperados: Adaptabilidad", [
        "Espectro de control en robótica:",
        "Teleoperación completa:",
        ("Humano controla cada movimiento", 1),
        ("Alta precisión, requiere operador constante", 1),
        "Teleoperación supervisada:",
        ("Robot ejecuta tareas, humano interviene cuando necesario", 1),
        ("Balance eficiencia-control", 1),
        "Autonomía guiada:",
        ("Objetivos de alto nivel por humano", 1),
        ("Robot planifica y ejecuta detalles", 1),
        "Autonomía completa:",
        ("Robot opera independientemente", 1),
        ("Humano solo recibe reportes", 1),
        "Tendencia: Transiciones dinámicas entre modos según contexto",
        "Aplicación: cirugía, manufactura, exploración espacial"
    ])

    # SECCIÓN 10: Adaptación
    create_section_slide(prs, "Adaptación y Control de Autonomía")

    # SLIDE 38: Adaptando Autonomía LLM
    create_content_slide(prs, "Adaptando Autonomía de Sistemas Basados en LLMs", [
        "Parámetros ajustables de autonomía:",
        "Temperatura y sampling:",
        ("Alta temperatura: más creatividad/riesgo", 1),
        ("Baja temperatura: más determinístico/seguro", 1),
        "Umbrales de confianza:",
        ("Auto-ejecuta si confianza > umbral", 1),
        ("Pide confirmación si confianza < umbral", 1),
        "Tool access permissions:",
        ("Lista blanca de herramientas permitidas", 1),
        ("Operaciones críticas requieren aprobación", 1),
        "Budget limits:",
        ("Tokens máximos por tarea", 1),
        ("Costo monetario máximo de API calls", 1),
        ("Tiempo máximo de ejecución", 1),
        "Configuración por contexto: desarrollo vs producción vs crítico"
    ])

    # SLIDE 39: Estados y Modos
    create_content_slide(prs, "Estados y Modos de Operación Autónomos", [
        "Máquinas de estado para control de agentes:",
        "Modo Observación:",
        ("Solo monitorea, no actúa", 1),
        ("Aprendizaje de patrones", 1),
        "Modo Sugerencia:",
        ("Propone acciones, espera aprobación", 1),
        ("Humano tiene control final", 1),
        "Modo Semi-Autónomo:",
        ("Ejecuta tareas rutinarias", 1),
        ("Escala decisiones críticas", 1),
        "Modo Autónomo:",
        ("Opera independientemente", 1),
        ("Reporta resultados periódicamente", 1),
        "Transiciones dinámicas:",
        ("Detección de anomalías → escalamiento automático", 1),
        ("Aprendizaje progresivo: ganando autonomía con confianza probada", 1)
    ])

    # SECCIÓN 11: Futuro
    create_section_slide(prs, "Visión del Futuro de IA Autónoma")

    # SLIDE 40: Futuro Cercano
    create_content_slide(prs, "Futuro Cercano (2025-2027): Agentes Productivos", [
        "2025 - Presente:",
        ("Adopción masiva de agentes en desarrollo de software", 1),
        ("25% empresas con pilotos de agentic AI", 1),
        ("Coding assistants como estándar en IDEs", 1),
        "2026:",
        ("Sistemas que descubren insights novedosos", 1),
        ("Agentes de larga duración (días a semanas) en producción", 1),
        ("Early AGI-like systems emergen (según proyecciones)", 1),
        "2027:",
        ("Robots autónomos en tareas del mundo real (manufactura, logística)", 1),
        ("50% de empresas usando agentic AI", 1),
        ("Primeros equipos 100% agentes en áreas específicas", 1),
        "Énfasis: Transición de pruebas de concepto a operaciones críticas",
        "Estudiantes 2025: Están en el momento perfecto para dominar esto"
    ])

    # SLIDE 41: Futuro Intermedio
    create_content_slide(prs, "Futuro Intermedio (2028-2030): AGI", [
        "Predicciones AGI (Artificial General Intelligence):",
        "Líderes industria (optimistas):",
        ("Sam Altman: AGI en 2026, superinteligencia 2030", 1),
        ("Dario Amodei: singularity 2026", 1),
        ("Eric Schmidt: AGI en 3-5 años (desde abril 2025)", 1),
        "Investigadores IA (conservadores):",
        ("Mediana: 50% probabilidad AGI en 2047", 1),
        ("90% probabilidad antes de 2075", 1),
        "Compresión acelerada de timelines:",
        ("Estimados cayeron de 50 años a 5 años en últimos 4 años", 1),
        "50% probabilidad varios hitos generales para 2028",
        "Implicaciones 2028-2030:",
        ("Sistemas que aprenden cualquier tarea intelectual humana", 1),
        ("Revolución en educación, medicina, investigación", 1),
        ("Desafíos masivos de alignment y seguridad", 1)
    ])

    # SLIDE 42: Futuro Lejano
    create_content_slide(prs, "Futuro Lejano (2030+): Escenarios", [
        "Superinteligencia:",
        ("IA superando capacidad humana en todos los dominios", 1),
        ("Timeframe post-AGI: 2-30 años según expertos", 1),
        "Escenarios transformacionales:",
        "Optimista:",
        ("Solución a cambio climático, enfermedades, pobreza", 1),
        ("Abundancia material y conocimiento ilimitado", 1),
        ("Colaboración humano-IA en nuevas fronteras", 1),
        "Cauteloso:",
        ("Disrupciones masivas de mercados laborales", 1),
        ("Necesidad de frameworks de governance global", 1),
        ("Desafíos de desigualdad de acceso a IA avanzada", 1),
        "Incertidumbre fundamental: Ningún experto puede predecir con certeza",
        "Responsabilidad: Generación actual diseñará estas tecnologías"
    ])

    # SECCIÓN 12: Práctica
    create_section_slide(prs, "Guía Práctica para Estudiantes")

    # SLIDE 43: Stack para Empezar
    create_content_slide(prs, "Herramientas Actuales: Stack para Empezar Hoy", [
        "LLMs (empezar gratis/económico):",
        ("OpenAI API (GPT-4o-mini), Anthropic (Claude)", 1),
        ("Open source local: Llama 3.3 8B, Qwen 2.5", 1),
        ("Plataformas: Ollama para local, OpenRouter para acceso múltiple", 1),
        "Frameworks de agentes:",
        ("Beginner: CrewAI - documentación excelente, rápido", 1),
        ("Intermediate: LangChain/LangGraph - ecosystem completo", 1),
        ("Advanced: AutoGen - multi-agent conversations", 1),
        "Desarrollo:",
        ("IDE: Cursor (student license), VS Code + Cline", 1),
        ("Deployment: Replit, Vercel, Railway", 1),
        "Automatización:",
        ("n8n (self-hosted gratis) - workflows con IA", 1),
        "Aprendizaje: GitHub repos, tutoriales YouTube, comunidades Discord"
    ])

    # SLIDE 44: Proyectos Prototipado
    create_content_slide(prs, "Proyectos de Prototipado: Ideas y Recursos", [
        "Proyectos iniciales (1-2 semanas):",
        "Personal AI assistant básico:",
        ("Gestión de tareas + integración calendario", 1),
        ("Stack: CrewAI + n8n + Google Calendar API", 1),
        "Research agent:",
        ("Busca papers, resume findings, genera reporte", 1),
        ("Stack: LangChain + RAG + Arxiv API", 1),
        "Code reviewer autónomo:",
        ("Analiza PRs, sugiere mejoras, detecta bugs", 1),
        ("Stack: AutoGen + GitHub API", 1),
        "Proyectos intermedios (1 mes):",
        ("Multi-agent customer support system", 1),
        ("Autonomous data analysis pipeline", 1),
        ("Content creation team (research + write + edit)", 1),
        "Recursos: Build in public, documentar en GitHub, compartir aprendizajes"
    ])

    # SLIDE 45: Deployment Real
    create_content_slide(prs, "Deployment Real: De la Idea a Producción", [
        "Ruta recomendada para estudiantes:",
        "1. Prototipo local (días):",
        ("Validar idea con Jupyter notebooks o scripts", 1),
        ("Iterar rápido sin preocuparse por infraestructura", 1),
        "2. MVP con usuarios (1-2 semanas):",
        ("Deploy en Replit/Vercel para feedback real", 1),
        ("Implementar analytics básico", 1),
        "3. Escalamiento (1-2 meses):",
        ("Migrar a infraestructura robusta si hay tracción", 1),
        ("Añadir monitoring, error handling, rate limiting", 1),
        "Consideraciones críticas:",
        ("Costos: establecer budgets de API calls", 1),
        ("Seguridad: validar inputs, proteger secrets", 1),
        ("Ethics: uso responsable, transparencia con usuarios", 1),
        "¡Empiecen HOY! El mejor momento para construir con IA autónoma es AHORA"
    ])

    print("Generando presentación PPTX...")
    prs.save('/home/user/experiments2/Inteligencia_Artificial_Autonoma_2025.pptx')
    print("✅ Presentación generada exitosamente: Inteligencia_Artificial_Autonoma_2025.pptx")
    print(f"📊 Total de slides: {len(prs.slides)}")

if __name__ == "__main__":
    main()
